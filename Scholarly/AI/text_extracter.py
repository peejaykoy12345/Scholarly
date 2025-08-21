import os
from fitz import open as fitz_open
from pptx import Presentation
from easyocr import Reader

reader = Reader(['en'], gpu=False)  

supported_image_files = {'png', 'jpg', 'jpeg'}

def is_supported_image_file(path: str) -> bool:
    ext = os.path.splitext(path)[1].lower().lstrip('.')
    return ext in supported_image_files

def extract_handwritten_text(image_path: str) -> str:
    result = reader.readtext(image_path, detail=0)
    return "\n".join(result) if result else "No text found."

def extract_pdf_text(pdf_path: str) -> str:
    doc = fitz_open(pdf_path)
    text = ""
    for page in doc:    
        text += page.get_text()
        #print(text)
    #print(f"Text: {text}")
    return text.strip() if text else "No text found."

def extract_pptx_text(pptx_path: str) -> str:
    prs = Presentation(pptx_path)

    all_text = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                all_text.append(shape.text)

    return "\n".join(all_text)

def extract_text(path: str) -> str:
    if path.lower().endswith('.pdf'):
        return extract_pdf_text(path)
    elif path.lower().endswith('.pptx'):
        return extract_pptx_text(path)
    elif is_supported_image_file(path):
        return extract_handwritten_text(path)
    else:
        return "❌ Unsupported file type."
